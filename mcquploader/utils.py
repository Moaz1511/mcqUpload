import re
import pandas as pd
import gspread
from pptx import Presentation
from oauth2client.service_account import ServiceAccountCredentials
from google.oauth2.service_account import Credentials
import tempfile
import shutil
from django.http import HttpResponse

def extract_mcq_info(text):
    # Define the regular expression pattern to match questions, references, options, answers, and explanations
    question_regex = r"(\d+।)\s*(.*?)\s*(?:\[(.*?)\])?\s+\(ক\)\s*(.*?)\s+\(খ\)\s*(.*?)\s+\(গ\)\s*(.*?)\s+\(ঘ\)\s*(.*?)\s+উত্তর:\s+(.*?)(?:\s+ব্যাখ্যা:\s+(.*?))?(?=\d+।|$)"
    
    match = re.finditer(question_regex, text, re.DOTALL)
    mcq_list = []
    for m in match:
        mcq_list.append(m.groups())
    return mcq_list

def process_pptx(file_path):
    # Load PowerPoint presentation
    presentation = Presentation(file_path)

    # Authenticate and open Google Sheets
    scope = ["https://spreadsheets.google.com/feeds", "https://www.googleapis.com/auth/drive"]
    
    credentials_path = "I:/10/Work From Home/13-03-2024/Website/Django/TwigTech/mcquploader/credentials.json"
    credentials = ServiceAccountCredentials.from_json_keyfile_name(credentials_path, scope)

    client = gspread.authorize(credentials)
    sheet = client.open("Your Google Sheet").sheet1

    row = 2  # Start from row 2 to avoid header

    batch_updates = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                mcq_info_list = extract_mcq_info(text)
                for mcq_info in mcq_info_list:
                    # Define the mapping of columns in the Google Sheet
                    # Change the column numbers as needed
                    serial_number_col = 1  # Column for serial number
                    question_col = 2  # Column for questions
                    reference_col = 3  # Column for reference
                    option_k_col = 4  # Column for option (ক)
                    option_kh_col = 5  # Column for option (খ)
                    option_g_col = 6  # Column for option (গ)
                    option_gh_col = 7  # Column for option (ঘ)
                    answer_col = 8  # Column for answers
                    explanation_col = 9  # Column for ব্যাখ্যা

                    # Initialize variables for each column value
                    serial_number, question, reference, option_k, option_kh, option_g, option_gh, answer, explanation = "", "", "", "", "", "", "", "", ""

                    # Extract all available data
                    if mcq_info:
                        serial_number, question, reference, option_k, option_kh, option_g, option_gh, answer, explanation = mcq_info

                    # Check if answer and explanation are empty and handle them
                    if not answer:
                        answer = ""
                    if not explanation:
                        explanation = ""

                    # Append the update as a list of values with columns
                    update_values = [
                        serial_number,
                        question,
                        reference,
                        option_k,
                        option_kh,
                        option_g,
                        option_gh,
                        answer,
                        explanation,
                    ]

                    batch_updates.append(update_values)

    # Update the Google Sheet in batches
    if batch_updates:
        sheet.update(f'A{row}:I{row + len(batch_updates) - 1}', batch_updates)

# Remember to replace 'path/to/credentials.json' with the actual path to your Google service account credentials
# and "Your Google Sheet" with the name of your Google Sheet.
'''
def export_worksheet_as_excel(spreadsheet_id, worksheet_title):
    # Authenticate and open Google Sheets
    # Define the scope
    scope = ["https://spreadsheets.google.com/feeds", "https://www.googleapis.com/auth/spreadsheets", "https://www.googleapis.com/auth/drive"]

    # Path to your service account credentials .json file
    credentials_path = "I:/10/Work From Home/13-03-2024/Website/Django/TwigTech/mcquploader/credentials.json"

    # Load the credentials
    credentials = Credentials.from_service_account_file(credentials_path, scopes=scope)

    # Authorize the client with gspread
    client = gspread.authorize(credentials)

    # Now you can use the client to open a Google Sheet by title or key
    sheet = client.open("Your Google Sheet").sheet1

    try:
        # Attempt to open the spreadsheet
        spreadsheet = client.open_by_key(spreadsheet_id)
        worksheet = spreadsheet.worksheet(worksheet_title)
        
        # Export the worksheet to a new temporary Excel file
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmpfile:
            worksheet.export(file_format='xlsx', filename=tmpfile.name)
            
            # Read the content of the temporary Excel file
            with open(tmpfile.name, 'rb') as f:
                excel_data = f.read()
            
            # Prepare the HttpResponse with the Excel file
            response = HttpResponse(excel_data, content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
            response['Content-Disposition'] = f'attachment; filename="{worksheet_title}.xlsx"'
            
            return response
    finally:
        # Cleanup the temporary file
        if 'tmpfile' in locals():
            shutil.rmtree(tmpfile.name, ignore_errors=True)
'''

def export_worksheet_as_excel(spreadsheet_id, worksheet_title):
    # Authenticate with Google Sheets
    credentials = Credentials.from_service_account_file('I:/10/Work From Home/13-03-2024/Website/Django/TwigTech/mcquploader/credentials.json', scopes=["https://www.googleapis.com/auth/spreadsheets"])
    gc = gspread.authorize(credentials)

    # Open the spreadsheet and worksheet
    sh = gc.open_by_key(spreadsheet_id)
    worksheet = sh.worksheet(worksheet_title)

    # Get all values in the worksheet
    data = worksheet.get_all_values()
    # Convert to a pandas DataFrame
    df = pd.DataFrame(data)
    df.columns = df.iloc[0] # Set first row as column names
    df = df.iloc[1:] # Remove first row

    # Create a Pandas Excel writer using XlsxWriter as the engine
    excel_file = f"{worksheet_title}.xlsx"
    sheet_name = 'Sheet1'
    
    with pd.ExcelWriter(excel_file, engine='openpyxl') as writer:
        df.to_excel(writer, sheet_name=sheet_name, index=False)

    return excel_file

